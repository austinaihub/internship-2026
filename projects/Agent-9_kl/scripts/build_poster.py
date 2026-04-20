"""
Build Kaiwen Lin's Capstone Poster — 36" x 36", UT Orange + black/white.
Output: outputs/12-Lin-Sp26.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from PIL import Image, ImageChops, ImageDraw


# ── Design tokens ────────────────────────────────────────────────────────────
ORANGE       = RGBColor(0xBF, 0x57, 0x00)
ORANGE_SOFT  = RGBColor(0xF7, 0xE6, 0xD6)
BLACK        = RGBColor(0x1D, 0x1D, 0x1F)
GRAY_DARK    = RGBColor(0x6E, 0x6E, 0x73)
GRAY_MID     = RGBColor(0xAE, 0xAE, 0xB2)
GRAY_LIGHT   = RGBColor(0xD2, 0xD2, 0xD7)
GRAY_BG      = RGBColor(0xFB, 0xFB, 0xFD)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW       = RGBColor(0xFF, 0xC7, 0x00)

FONT_HEAD    = "Helvetica Neue"
FONT_BODY    = "Helvetica Neue"

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets" / "logos"
OUT_DIR = ROOT / "outputs"
OUT_DIR.mkdir(parents=True, exist_ok=True)
STAGE_DIR = OUT_DIR / "_logo_stage"
STAGE_DIR.mkdir(parents=True, exist_ok=True)
REVIEWER_PREVIEW = OUT_DIR / "reviewer-website-preview.png"


# ── Image utils ──────────────────────────────────────────────────────────────
def _trim_uniform_bg(img: Image.Image, bg_rgb: tuple, tol: int = 12) -> Image.Image:
    """Trim a uniform-colored border. bg_rgb is the color to trim."""
    rgb = img.convert("RGB")
    ref = Image.new("RGB", rgb.size, bg_rgb)
    diff = ImageChops.difference(rgb, ref)
    # Threshold: pixels within `tol` of bg_rgb → treated as bg
    bw = diff.point(lambda p: 255 if p > tol else 0)
    bbox = bw.getbbox()
    return img.crop(bbox) if bbox else img


def prepare_ischool(src: Path) -> Path:
    """Trim white border from iSchool logo."""
    dst = STAGE_DIR / "ischool_prep.png"
    img = Image.open(src).convert("RGBA")
    trimmed = _trim_uniform_bg(img, (255, 255, 255), tol=10)
    trimmed.save(dst)
    return dst


def prepare_austinaihub(src: Path) -> Path:
    """Trim dark-navy padding and add rounded corners for a 'card' feel."""
    dst = STAGE_DIR / "austinaihub_prep.png"
    img = Image.open(src).convert("RGBA")
    # Sample corner to detect bg color
    corner = img.getpixel((1, 1))[:3]
    trimmed = _trim_uniform_bg(img, corner, tol=18)
    # Add small rounded-corner mask for a clean sponsor-card look
    w, h = trimmed.size
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    radius = max(10, min(w, h) // 14)
    draw.rounded_rectangle([(0, 0), (w, h)], radius=radius, fill=255)
    out = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    out.paste(trimmed, (0, 0), mask)
    out.save(dst)
    return dst


# ── PPT helpers ──────────────────────────────────────────────────────────────
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width_pt=1.0):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, *, size=22, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=WHITE, line_color=GRAY_LIGHT,
             line_w=1.0, rounded=True, no_outline=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if rounded:
        try:
            s.adjustments[0] = 0.06
        except Exception:
            pass
    set_fill(s, fill)
    if no_outline:
        no_line(s)
    else:
        set_line(s, line_color, line_w)
    s.text_frame.text = ""
    return s


def add_card(slide, x, y, w, h, title, bullets, *, accent=ORANGE, body_size=22):
    """A titled card: orange accent bar, section header, and bullets."""
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x + 0.35), Inches(y + 0.35),
                                 Inches(0.35), Inches(0.9))
    set_fill(bar, accent); no_line(bar)
    add_text(slide, x + 0.9, y + 0.3, w - 1.1, 1.0, title,
             size=40, bold=True, color=BLACK, font=FONT_HEAD)
    tb = slide.shapes.add_textbox(Inches(x + 0.5), Inches(y + 1.5),
                                  Inches(w - 1.0), Inches(h - 1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        if " — " in b:
            head, rest = b.split(" — ", 1)
            r1 = p.add_run()
            r1.text = f"•  {head} — "
            r1.font.name = FONT_BODY; r1.font.size = Pt(body_size); r1.font.bold = True
            r1.font.color.rgb = BLACK
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = FONT_BODY; r2.font.size = Pt(body_size); r2.font.bold = False
            r2.font.color.rgb = GRAY_DARK
        else:
            r = p.add_run()
            r.text = f"•  {b}"
            r.font.name = FONT_BODY; r.font.size = Pt(body_size); r.font.color.rgb = GRAY_DARK


def add_placeholder(slide, x, y, w, h, label, *, accent=ORANGE):
    add_rect(slide, x, y, w, h, fill=GRAY_BG, line_color=accent,
             line_w=2.0, rounded=True)
    add_text(slide, x, y + h/2 - 0.3, w, 0.5, label,
             size=22, color=GRAY_DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, italic=True)


def add_stack_chip(slide, x, y, w, h, primary, secondary):
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=ORANGE,
             line_w=1.8, rounded=True)
    add_text(slide, x, y + 0.12, w, 0.28, primary,
             size=16, bold=True, color=ORANGE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, x, y + 0.42, w, 0.20, secondary,
             size=10, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.CENTER, italic=True, wrap=False)


def add_ui_mockup(slide, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=ORANGE,
             line_w=2.0, rounded=True)

    shell_x, shell_y = x + 0.18, y + 0.18
    shell_w, shell_h = w - 0.36, h - 0.36
    add_rect(slide, shell_x, shell_y, shell_w, shell_h, fill=WHITE,
             line_color=GRAY_LIGHT, rounded=True)

    bar_h = 0.46
    add_rect(slide, shell_x, shell_y, shell_w, bar_h, fill=GRAY_BG,
             line_color=GRAY_LIGHT, rounded=True)
    dot_y = shell_y + 0.14
    for i, color in enumerate([ORANGE, GRAY_MID, GRAY_MID]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(shell_x + 0.14 + i * 0.16), Inches(dot_y),
            Inches(0.09), Inches(0.09))
        set_fill(dot, color)
        no_line(dot)
    add_text(slide, shell_x + 0.58, shell_y + 0.10, 2.5, 0.16,
             "Campaign Agent reviewer", size=10, bold=True,
             color=BLACK, font=FONT_BODY, wrap=False)
    add_text(slide, shell_x + shell_w - 2.0, shell_y + 0.10, 1.7, 0.16,
             "session paused", size=10, color=GRAY_DARK,
             font=FONT_BODY, align=PP_ALIGN.RIGHT, wrap=False)

    content_y = shell_y + bar_h + 0.12
    content_h = shell_h - bar_h - 0.24
    side_w = 1.65
    gutter = 0.20
    list_w = 3.15
    preview_w = shell_w - side_w - list_w - 2 * gutter - 0.24

    add_rect(slide, shell_x + 0.12, content_y, side_w, content_h,
             fill=GRAY_BG, line_color=GRAY_LIGHT, rounded=True)
    add_text(slide, shell_x + 0.26, content_y + 0.14, side_w - 0.28, 0.22,
             "Review lanes", size=11, bold=True, color=BLACK,
             font=FONT_HEAD, wrap=False)
    sidebar_chips = [
        ("TREND", ORANGE_SOFT, ORANGE),
        ("NEEDS REVIEW", WHITE, BLACK),
        ("CONTENT", WHITE, BLACK),
        ("IMAGE", WHITE, BLACK),
    ]
    for i, (label, fill, text_color) in enumerate(sidebar_chips):
        chip_y = content_y + 0.48 + i * 0.56
        add_rect(slide, shell_x + 0.24, chip_y, side_w - 0.24, 0.36,
                 fill=fill, line_color=GRAY_LIGHT, rounded=True)
        add_text(slide, shell_x + 0.24, chip_y + 0.06, side_w - 0.24, 0.14,
                 label, size=9, bold=True, color=text_color,
                 font=FONT_BODY, align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, shell_x + 0.26, content_y + content_h - 0.74,
             side_w - 0.30, 0.44,
             "3 gates only.\nHuman checks the sensitive moments.",
             size=9, color=GRAY_DARK, font=FONT_BODY,
             italic=True, wrap=False)

    list_x = shell_x + side_w + gutter
    add_rect(slide, list_x, content_y, list_w, content_h,
             fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    add_text(slide, list_x + 0.16, content_y + 0.14, list_w - 0.32, 0.22,
             "Trend queue", size=11, bold=True, color=BLACK,
             font=FONT_HEAD, wrap=False)
    add_text(slide, list_x + list_w - 1.2, content_y + 0.14, 1.0, 0.22,
             "5 articles", size=9, color=GRAY_DARK,
             font=FONT_BODY, align=PP_ALIGN.RIGHT, wrap=False)

    article_cards = [
        ("South Korea addresses child trafficking", "Approved topic"),
        ("Florida civil liberties debate gains steam", "Selected now"),
        ("Local hotline funding falls behind demand", "Queued"),
    ]
    for i, (headline, state) in enumerate(article_cards):
        card_y = content_y + 0.45 + i * 0.86
        border = ORANGE if i == 1 else GRAY_LIGHT
        fill = ORANGE_SOFT if i == 1 else WHITE
        add_rect(slide, list_x + 0.14, card_y, list_w - 0.28, 0.68,
                 fill=fill, line_color=border, line_w=1.8, rounded=True)
        if i == 1:
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(list_x + 0.14), Inches(card_y),
                Inches(0.08), Inches(0.68))
            set_fill(accent, ORANGE)
            no_line(accent)
        add_text(slide, list_x + 0.28, card_y + 0.10, list_w - 0.52, 0.22,
                 headline, size=9, bold=True, color=BLACK,
                 font=FONT_BODY, wrap=False)
        add_text(slide, list_x + 0.28, card_y + 0.36, 1.35, 0.16,
                 state, size=8, color=GRAY_DARK,
                 font=FONT_BODY, italic=True, wrap=False)
        add_text(slide, list_x + list_w - 1.00, card_y + 0.36, 0.72, 0.16,
                 "Open", size=8, bold=True, color=ORANGE,
                 font=FONT_BODY, align=PP_ALIGN.RIGHT, wrap=False)

    preview_x = list_x + list_w + gutter
    add_rect(slide, preview_x, content_y, preview_w, content_h,
             fill=GRAY_BG, line_color=GRAY_LIGHT, rounded=True)
    add_text(slide, preview_x + 0.18, content_y + 0.14, preview_w - 0.36, 0.22,
             "Draft preview", size=11, bold=True, color=BLACK,
             font=FONT_HEAD, wrap=False)
    add_pill(slide, preview_x + 0.18, content_y + 0.42, 1.35, 0.30,
             "AUDIENCE", fill=ORANGE_SOFT, text_color=ORANGE,
             size=9, bold=True)
    add_text(slide, preview_x + 1.70, content_y + 0.46, preview_w - 1.88, 0.16,
             "people who trust the wrong person", size=8,
             color=GRAY_DARK, font=FONT_BODY, wrap=False)
    add_text(slide, preview_x + 0.18, content_y + 0.86, preview_w - 0.36, 0.44,
             "The warning signs rarely look like danger at first.",
             size=12, bold=True, color=BLACK, font=FONT_HEAD, wrap=False)

    body_lines = [
        "Traffickers often arrive as someone familiar.",
        "A partner, recruiter, or friend can still be the threat.",
        "Early, human-reviewed messages matter before harm escalates.",
    ]
    for i, line in enumerate(body_lines):
        add_text(slide, preview_x + 0.18, content_y + 1.34 + i * 0.34,
                 preview_w - 0.36, 0.18, line, size=8.5,
                 color=BLACK if i == 0 else GRAY_DARK,
                 font=FONT_BODY, wrap=False)

    add_text(slide, preview_x + 0.18, content_y + content_h - 1.16,
             preview_w - 0.36, 0.18, "Decision controls",
             size=9, bold=True, color=BLACK, font=FONT_BODY, wrap=False)
    button_specs = [
        ("Pause", WHITE, ORANGE),
        ("Approve", ORANGE, WHITE),
        ("Refine", WHITE, ORANGE),
    ]
    btn_w = (preview_w - 0.48 - 0.20) / 3
    for i, (label, fill, text_color) in enumerate(button_specs):
        btn_x = preview_x + 0.18 + i * (btn_w + 0.10)
        btn_y = content_y + content_h - 0.76
        add_rect(slide, btn_x, btn_y, btn_w, 0.42, fill=fill,
                 line_color=ORANGE, line_w=1.4, rounded=True)
        add_text(slide, btn_x, btn_y + 0.08, btn_w, 0.14, label,
                 size=9, bold=True, color=text_color,
                 font=FONT_BODY, align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, preview_x + 0.18, content_y + content_h - 0.22,
             preview_w - 0.36, 0.16,
             "Reviewer approves, edits, or sends the draft back for another pass.",
             size=8.5, color=GRAY_DARK, font=FONT_BODY,
             italic=True, wrap=False)


def add_narrative_card(slide, x, y, w, h, title, pull_quote, paragraphs,
                       closing, *, body_size=18, closing_size=20):
    """A narrative card: title + orange pull quote + body paragraphs + orange closing line."""
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x + 0.35), Inches(y + 0.35),
                                 Inches(0.35), Inches(0.9))
    set_fill(bar, ORANGE); no_line(bar)
    add_text(slide, x + 0.9, y + 0.3, w - 1.1, 1.0, title,
             size=40, bold=True, color=BLACK, font=FONT_HEAD)

    # Pull quote
    add_text(slide, x + 0.5, y + 1.45, w - 1.0, 1.0, pull_quote,
             size=26, bold=True, color=ORANGE, font=FONT_HEAD, italic=True)

    # Body paragraphs
    tb = slide.shapes.add_textbox(Inches(x + 0.5), Inches(y + 2.55),
                                  Inches(w - 1.0), Inches(h - 4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = para
        r.font.name = FONT_BODY; r.font.size = Pt(body_size); r.font.color.rgb = BLACK

    # Closing — orange, bold, bottom of card
    add_text(slide, x + 0.5, y + h - 1.4, w - 1.0, 1.1, closing,
             size=closing_size, bold=True, color=ORANGE, font=FONT_BODY)


def add_numbered_step(slide, x, y, w, h, num, title, tech, desc, *,
                      hitl_label=None):
    """A horizontal step row: [number circle]  TITLE / tech / description  [HITL badge?]"""
    # Number circle (left)
    circle_d = 0.9
    circle_y = y + (h - circle_d) / 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x + 0.05), Inches(circle_y),
                                    Inches(circle_d), Inches(circle_d))
    set_fill(circle, ORANGE); no_line(circle)
    tf = circle.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.name = FONT_HEAD; r.font.size = Pt(32); r.font.bold = True
    r.font.color.rgb = WHITE

    # Content column (title + tech + desc)
    content_x = x + 1.1
    content_w = w - 1.2 - (2.5 if hitl_label else 0.0)
    add_text(slide, content_x, y + 0.03, content_w, 0.42, title,
             size=20, bold=True, color=BLACK, font=FONT_HEAD)
    add_text(slide, content_x, y + 0.44, content_w, 0.28, tech,
             size=13, color=ORANGE, font=FONT_BODY, italic=True)
    add_text(slide, content_x, y + 0.72, content_w, h - 0.75, desc,
             size=13, color=GRAY_DARK, font=FONT_BODY)

    # HITL badge (right, if present)
    if hitl_label:
        bw, bh = 2.3, 0.55
        bx = x + w - bw - 0.05
        by = y + (h - bh) / 2
        add_pill(slide, bx, by, bw, bh, hitl_label,
                 fill=YELLOW, text_color=BLACK, size=11, bold=True)


def add_pill_two_line(slide, x, y, w, h, primary, secondary, *,
                      fill=ORANGE, primary_color=WHITE, secondary_color=WHITE,
                      primary_size=20, secondary_size=12):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        s.adjustments[0] = 0.30
    except Exception:
        pass
    set_fill(s, fill); no_line(s)
    tf = s.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(2)
    r1 = p1.add_run(); r1.text = primary
    r1.font.name = FONT_HEAD; r1.font.size = Pt(primary_size); r1.font.bold = True
    r1.font.color.rgb = primary_color

    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = secondary
    r2.font.name = FONT_BODY; r2.font.size = Pt(secondary_size); r2.font.bold = False
    r2.font.italic = True
    r2.font.color.rgb = secondary_color


def insert_logo(slide, src_path: Path, bbox_x, bbox_y, bbox_w, bbox_h,
                *, placeholder_label: str):
    """Insert a logo auto-fit into a bbox, preserving aspect. Fallback to placeholder."""
    if not src_path or not Path(src_path).exists():
        add_placeholder(slide, bbox_x, bbox_y, bbox_w, bbox_h, placeholder_label)
        return
    im = Image.open(src_path)
    src_w, src_h = im.size
    src_ratio = src_w / src_h
    box_ratio = bbox_w / bbox_h
    if src_ratio > box_ratio:
        render_w = bbox_w
        render_h = bbox_w / src_ratio
    else:
        render_h = bbox_h
        render_w = bbox_h * src_ratio
    px = bbox_x + (bbox_w - render_w) / 2
    py = bbox_y + (bbox_h - render_h) / 2
    slide.shapes.add_picture(str(src_path), Inches(px), Inches(py),
                             Inches(render_w), Inches(render_h))


def add_pill(slide, x, y, w, h, text, *, fill=ORANGE_SOFT, text_color=ORANGE,
             size=18, bold=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        s.adjustments[0] = 0.5
    except Exception:
        pass
    set_fill(s, fill); no_line(s)
    tf = s.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_HEAD; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = text_color


def add_connector_arrow(slide, x1, y1, x2, y2, *, color=ORANGE, width_pt=3.0,
                        dashed=False):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    if dashed:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')


# ── Build the poster ─────────────────────────────────────────────────────────
def build():
    # Prepare logos
    ischool_src = ASSETS / "ischool.png"
    austin_src  = ASSETS / "austinaihub.png"
    ischool_prep = prepare_ischool(ischool_src) if ischool_src.exists() else None
    austin_prep  = prepare_austinaihub(austin_src) if austin_src.exists() else None

    prs = Presentation()
    prs.slide_width = Inches(36)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    W, H = 36.0, 36.0

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    set_fill(bg, GRAY_BG); no_line(bg)

    # ── HEADER (0 – 4.8") ────────────────────────────────────────────────────
    # iSchool logo (top-left) — tighter footprint so the header reads as one unit
    insert_logo(slide, ischool_prep, 1.0, 0.9, 4.7, 1.55,
                placeholder_label="iSchool logo — drop here")

    # Poster number (top-right) — keep it within the required 60–72 pt range
    add_text(slide, W - 4.6, 0.82, 3.6, 1.40, "12",
             size=72, bold=True, color=ORANGE, font=FONT_HEAD,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)
    add_text(slide, W - 4.6, 2.05, 3.6, 0.32, "Poster Number",
             size=15, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

    # Title — compressed vertically so the page gets a clearer visual center
    title_x = 7.0
    title_w = 22.0
    add_text(slide, title_x, 0.62, title_w, 1.15,
             "Human-in-the-Loop Multi-Agent System",
             size=62, bold=True, color=BLACK, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(slide, title_x, 1.58, title_w, 0.85,
             "for Social Campaign Automation",
             size=44, bold=True, color=ORANGE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Byline (3 lines) — tucked closer to the title so the header feels intentional
    add_text(slide, 1.0, 2.62, W - 2.0, 0.46,
             "Kaiwen Lin  —  AI Engineer Intern @ AustinAIHUB  ·  UTEID  kl38264",
             size=24, bold=True, color=BLACK, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.0, 3.10, W - 2.0, 0.30,
             "Team: Amirhossein Azami  ·  Armin Mohammadi",
             size=17, color=BLACK, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.0, 3.50, W - 2.0, 0.28,
             "Supervisor: Ali Mirzapour   ·   UT iSchool Capstone — Spring 2026",
             size=14, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.CENTER)

    # Orange rule under header
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.20),
                                  Inches(W - 2.0), Inches(0.08))
    set_fill(rule, ORANGE); no_line(rule)

    # ── MIDDLE (5.1 – 15.3") — Abstract | Architecture | Methods ────────────
    Y = 4.65
    COL_H = 10.45
    LEFT_W = 9.0
    CENTER_X = 10.3
    CENTER_W = 15.4
    RIGHT_X = 26.0
    RIGHT_W = 9.0

    hidden_crisis_paragraphs = [
        "An anti-human-trafficking campaign is rarely speaking to someone who already says, \"I am being trafficked.\" It may be reaching a teenager who thinks an older boyfriend is helping, a runaway promised a couch and cash, or a worker told they only owe one more favor.",
        "That is what makes trafficking so hard to interrupt: coercion arrives disguised as love, debt, safety, rides, gifts, and shame. By the time the story becomes a headline, the manipulation has usually been rehearsed for weeks or months.",
        "So the real challenge is earlier. Anti-trafficking organizations need public messages that sound human, move fast, and meet people inside that confusion — before the person at risk would ever search the word \"trafficking.\"",
    ]
    add_narrative_card(
        slide, 1.0, Y, LEFT_W, COL_H,
        title="Before It Looks Like Trafficking",
        pull_quote="The person most at risk may not call it trafficking yet.",
        paragraphs=hidden_crisis_paragraphs,
        closing="Campaign Agent turns breaking news into audience-tuned anti-trafficking warnings in minutes, with a human reviewing every sensitive decision.",
        body_size=17,
        closing_size=19,
    )

    # Architecture center
    add_rect(slide, CENTER_X, Y, CENTER_W, COL_H,
             fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(CENTER_X + 0.35), Inches(Y + 0.35),
                                 Inches(0.35), Inches(0.9))
    set_fill(bar, ORANGE); no_line(bar)
    add_text(slide, CENTER_X + 0.9, Y + 0.3, CENTER_W - 1.1, 1.0,
             "What We Built",
             size=40, bold=True, color=BLACK, font=FONT_HEAD)
    add_text(slide, CENTER_X + 0.9, Y + 1.15, CENTER_W - 1.1, 0.5,
             "Five specialist agents drafting.  One human deciding.",
             size=22, bold=True, color=GRAY_DARK, italic=True)

    # Supervisor
    sup_x, sup_y, sup_w, sup_h = CENTER_X + 5.2, Y + 1.9, 5.0, 1.1
    sup = add_rect(slide, sup_x, sup_y, sup_w, sup_h,
                   fill=ORANGE, line_color=ORANGE, rounded=True, no_outline=True)
    add_text(slide, sup_x, sup_y + 0.1, sup_w, sup_h - 0.1,
             ["Supervisor", "LangGraph router"],
             size=22, bold=True, color=WHITE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    agents = [
        ("Trend",     "Exa.ai",    "finds fresh stories"),
        ("Audience",  "ChatGPT",   "defines the voice"),
        ("Writer",    "ChatGPT",   "drafts the post"),
        ("Image",     "Gemini",    "renders the visual"),
        ("Publisher", "Tweepy",    "ships to social"),
    ]
    node_y = Y + 4.4
    node_w, node_h = 2.5, 1.65
    gap = (CENTER_W - 2 * 0.9 - 5 * node_w) / 4
    start_x = CENTER_X + 0.9
    node_centers = []
    for i, (name, svc, role) in enumerate(agents):
        nx = start_x + i * (node_w + gap)
        ny = node_y
        add_rect(slide, nx, ny, node_w, node_h, fill=WHITE,
                 line_color=ORANGE, line_w=2.5, rounded=True)
        # Name (22pt bold)
        add_text(slide, nx, ny + 0.08, node_w, 0.42, name,
                 size=22, bold=True, color=BLACK, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER)
        # Service (16pt)
        add_text(slide, nx, ny + 0.54, node_w, 0.32, svc,
                 size=16, color=GRAY_DARK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER)
        # Role line (13pt italic — readable from 1.5m viewing distance)
        add_text(slide, nx + 0.05, ny + 0.90, node_w - 0.1, 0.70, role,
                 size=13, italic=True, color=GRAY_DARK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER)
        node_centers.append((nx + node_w/2, ny, ny + node_h))

    for cx, top_y, _ in node_centers:
        add_connector_arrow(slide, sup_x + sup_w/2, sup_y + sup_h,
                            cx, top_y, color=ORANGE, width_pt=2.5)

    mid_y = node_y + node_h/2
    flow_labels = ["topic", "brief", "draft", "caption"]
    for i in range(len(node_centers) - 1):
        x_from = node_centers[i][0] + node_w/2
        x_to   = node_centers[i+1][0] - node_w/2
        add_connector_arrow(slide, x_from, mid_y, x_to, mid_y,
                            color=GRAY_MID, width_pt=2.0)
        # Label above the arrow — 12pt bold black for maximum contrast
        lbl_cx = (x_from + x_to) / 2
        lbl_w = 1.0
        add_text(slide, lbl_cx - lbl_w/2, mid_y - 0.28, lbl_w, 0.26,
                 flow_labels[i],
                 size=12, bold=True, color=BLACK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER)

    hitl_indices = [0, 1, 3]
    hitl_labels = ["HITL: Topic", "HITL: Audience", "HITL: Image"]
    for idx, label in zip(hitl_indices, hitl_labels):
        cx, _, bot_y = node_centers[idx]
        by = bot_y + 0.35
        bw = 2.2; bh = 0.6
        add_pill(slide, cx - bw/2, by, bw, bh, label,
                 fill=YELLOW, text_color=BLACK)
        add_connector_arrow(slide, cx - 0.3, by, cx - 0.3, bot_y,
                            color=GRAY_DARK, width_pt=1.5, dashed=True)

    # ── Info panels: State Management (LangGraph) + Observability (LangSmith)
    panel_y = Y + 7.30   # y = 12.40
    panel_h = 1.60
    panel_gap = 0.25
    panel_w = (CENTER_W - 1.0 - panel_gap) / 2  # 7.075
    panel_left_x  = CENTER_X + 0.5
    panel_right_x = panel_left_x + panel_w + panel_gap

    def _info_panel(px, py, pw, ph, title, body):
        add_rect(slide, px, py, pw, ph,
                 fill=GRAY_BG, line_color=GRAY_LIGHT, rounded=True)
        # Left orange accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(px + 0.20), Inches(py + 0.20),
                                     Inches(0.22), Inches(0.55))
        set_fill(bar, ORANGE); no_line(bar)
        add_text(slide, px + 0.55, py + 0.15, pw - 0.75, 0.45, title,
                 size=17, bold=True, color=BLACK, font=FONT_HEAD)
        add_text(slide, px + 0.30, py + 0.75, pw - 0.55, ph - 0.85, body,
                 size=14, color=GRAY_DARK, font=FONT_BODY)

    _info_panel(panel_left_x, panel_y, panel_w, panel_h,
                "State Management · LangGraph",
                "A typed CampaignState flows agent-to-agent. interrupt() "
                "pauses the graph at THREE gates — topic, audience, image — "
                "not five. Judgment where it matters; automation where it doesn't.")
    _info_panel(panel_right_x, panel_y, panel_w, panel_h,
                "Observability · LangSmith",
                "Every prompt, retry, and token is traced per agent. "
                "Sensitive advocacy work demands receipts — LangSmith is "
                "the audit log for every decision the system makes.")

    add_text(slide, CENTER_X + 0.9, Y + COL_H - 0.8, CENTER_W - 1.8, 0.5,
             "React reviewer UI  ·  Pause → Approve → Resume  ·  Sensitive content never ships without a human",
             size=18, bold=True, color=ORANGE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)

    # ── RIGHT CARD: "Where Humans Win — HITL Authority System" ─────────────
    # Replaces bullet list with a 3-section visual diagram showing the
    # three HITL gates, the one-input-four-agents fan-out, and the five
    # priority-mechanism code patterns that make user intent structurally win.
    add_rect(slide, RIGHT_X, Y, RIGHT_W, COL_H,
             fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    bar_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(RIGHT_X + 0.35), Inches(Y + 0.35),
                                   Inches(0.35), Inches(0.9))
    set_fill(bar_r, ORANGE); no_line(bar_r)
    add_text(slide, RIGHT_X + 0.9, Y + 0.26, RIGHT_W - 1.1, 0.92,
             "Where Humans Win",
             size=34, bold=True, color=BLACK, font="Arial", wrap=False)
    add_text(slide, RIGHT_X + 0.9, Y + 1.18, RIGHT_W - 1.1, 0.28,
             "The HITL Authority System",
             size=14, italic=True, color=GRAY_DARK, font=FONT_BODY, wrap=False)

    card_inner_x = RIGHT_X + 0.4
    card_inner_w = RIGHT_W - 0.8

    # ── Section ①: Three gates, three decisions ─────────────────────────
    sec1_y = Y + 1.72
    add_text(slide, card_inner_x, sec1_y, card_inner_w, 0.32,
             "①  Three Gates, Three Decisions",
             size=15, bold=True, color=ORANGE, font=FONT_HEAD, wrap=False)

    # Vertical orange spine
    spine_x = card_inner_x + 0.30
    spine_top = sec1_y + 0.45
    spine_bot = sec1_y + 3.30
    spine = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(spine_x), Inches(spine_top),
        Inches(0.04), Inches(spine_bot - spine_top))
    set_fill(spine, ORANGE); no_line(spine)

    gates = [
        ("⏸  GATE 1  ·  Trend",
         "Approve  ·  Override topic  ·  Add guidance  ·  Re-search"),
        ("⏸  GATE 2  ·  Content",
         "Approve & publish  ·  Regen image  ·  Regen all + feedback"),
        ("⏸  GATE 3  ·  Refine",
         "Refine text / image / both / audience  (forks new session)"),
    ]
    gate_h = 0.55
    gate_gap = 0.40
    gate_x = spine_x + 0.25
    pill_w = 2.8
    for i, (pill_text, actions) in enumerate(gates):
        gy = spine_top + i * (gate_h + gate_gap + 0.25)
        add_pill(slide, gate_x, gy, pill_w, gate_h, pill_text,
                 fill=YELLOW, text_color=BLACK, size=14, bold=True)
        # Small orange dot on the spine at this gate
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(spine_x - 0.07), Inches(gy + gate_h/2 - 0.09),
            Inches(0.18), Inches(0.18))
        set_fill(dot, ORANGE); no_line(dot)
        # Action list below pill
        add_text(slide, gate_x, gy + gate_h + 0.04,
                 card_inner_w - (gate_x - card_inner_x), 0.30,
                 actions,
                 size=11, color=GRAY_DARK, font=FONT_BODY,
                 italic=True, wrap=False)

    # ── Section ②: One Input → Four Agents (fan-out) ───────────────────
    sec2_y = sec1_y + 3.55
    add_text(slide, card_inner_x, sec2_y, card_inner_w, 0.32,
             "②  One Input  →  Four Agents",
             size=15, bold=True, color=ORANGE, font=FONT_HEAD, wrap=False)

    # User input box at top center
    fan_center_x = RIGHT_X + RIGHT_W / 2
    user_w, user_h = 4.0, 0.70
    user_x = fan_center_x - user_w / 2
    user_y = sec2_y + 0.42
    user_box = add_rect(slide, user_x, user_y, user_w, user_h,
                        fill=ORANGE, line_color=ORANGE,
                        rounded=True, no_outline=True)
    add_text(slide, user_x, user_y + 0.05, user_w, user_h - 0.10,
             ["👤  user_guidance",
              "\"hopeful tone, survivor's strength\""],
             size=12, bold=True, color=WHITE, font=FONT_BODY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)

    # Three recipient boxes in a row below
    rec_y = user_y + user_h + 0.75
    rec_h = 0.55
    rec_w = 2.35
    rec_total = 3 * rec_w + 2 * 0.15
    rec_start_x = fan_center_x - rec_total / 2
    recipients = ["Audience", "Writer", "Image Gen"]
    rec_centers = []
    for i, name in enumerate(recipients):
        rx = rec_start_x + i * (rec_w + 0.15)
        add_rect(slide, rx, rec_y, rec_w, rec_h,
                 fill=WHITE, line_color=ORANGE, rounded=True)
        add_text(slide, rx, rec_y + 0.04, rec_w, rec_h - 0.08, name,
                 size=13, bold=True, color=BLACK, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        rec_centers.append(rx + rec_w / 2)

    # Fan-out arrows: user box bottom → each recipient top
    user_bot = user_y + user_h
    for cx in rec_centers:
        add_connector_arrow(slide, fan_center_x, user_bot, cx, rec_y,
                            color=ORANGE, width_pt=1.5)

    # Double-path: dashed arrow from Audience bottom → Image Gen bottom (curves under)
    aud_x = rec_centers[0]
    img_x = rec_centers[2]
    loop_y = rec_y + rec_h + 0.22
    add_connector_arrow(slide, aud_x, rec_y + rec_h, aud_x, loop_y,
                        color=ORANGE, width_pt=1.0, dashed=True)
    add_connector_arrow(slide, aud_x, loop_y, img_x, loop_y,
                        color=ORANGE, width_pt=1.0, dashed=True)
    add_connector_arrow(slide, img_x, loop_y, img_x, rec_y + rec_h,
                        color=ORANGE, width_pt=1.0, dashed=True)
    add_text(slide, card_inner_x, loop_y + 0.06, card_inner_w, 0.28,
             "indirect path: visual_style shapes the image prompt",
             size=10, italic=True, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, card_inner_x, loop_y + 0.38, card_inner_w, 0.28,
             "Image Gen reached via TWO paths — intent is structurally embedded",
             size=10, bold=True, color=BLACK, font=FONT_BODY,
             align=PP_ALIGN.CENTER, wrap=False)

    # ── Section ③: Five priority mechanisms ────────────────────────────
    sec3_y = sec2_y + 2.65
    add_text(slide, card_inner_x, sec3_y, card_inner_w, 0.32,
             "③  How User Intent Wins  —  5 Code Mechanisms",
             size=15, bold=True, color=ORANGE, font=FONT_HEAD, wrap=False)

    mechs = [
        ("Direct Override",       'update["trend_topic"] = req.custom_topic'),
        ("CRITICAL Label",        'prompt += "CRITICAL INSTRUCTION FROM USER..."'),
        ("Prompt Replacement",    'human_content = writer_prompt + ...'),
        ("Multi-Agent Propagate", 'user_guidance  →  3 agent prompts'),
        ("Selective Clear",       'post_text=None; keeps audience+image'),
    ]
    mech_y = sec3_y + 0.38
    mech_h = 0.36
    mech_gap = 0.06
    for i, (name, sig) in enumerate(mechs):
        my = mech_y + i * (mech_h + mech_gap)
        # Row background
        add_rect(slide, card_inner_x, my, card_inner_w, mech_h,
                 fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
        # Left: mechanism name (bold black)
        add_text(slide, card_inner_x + 0.15, my + 0.03,
                 2.55, mech_h - 0.06, name,
                 size=12, bold=True, color=BLACK, font=FONT_HEAD,
                 anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        # Right: code signature (mono-ish gray)
        add_text(slide, card_inner_x + 2.75, my + 0.03,
                 card_inner_w - 2.90, mech_h - 0.06, sig,
                 size=10, color=GRAY_DARK, font=FONT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE, wrap=False)

    # ── DEMO STRIP (15.4 – 23.2") ────────────────────────────────────────────
    # Three zones: (1) Website preview screenshot, (2) Real output images,
    # (3) How an Image Is Born — vertical 5-step flow
    demo_title_y = 15.10
    demo_subtitle_y = 15.65
    demo_y = 15.85
    demo_h = 7.50
    demo_w = 10.8
    demo_gap = (W - 2.0 - 3 * demo_w) / 2
    block_x = [1.0 + i * (demo_w + demo_gap) for i in range(3)]

    def _zone_title(x, y, w, text):
        add_text(slide, x, y, w, 0.5, text,
                 size=22, bold=True, color=BLACK, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER)
        # Small orange underline accent, centered
        bar_w = 3.0
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x + (w - bar_w) / 2),
                                     Inches(y + 0.5),
                                     Inches(bar_w), Inches(0.04))
        set_fill(bar, ORANGE); no_line(bar)

    # ── BLOCK 1 (left): Website preview screenshot ──────────────────────────
    _zone_title(block_x[0], demo_title_y, demo_w, "Website Preview")
    add_text(slide, block_x[0], demo_subtitle_y, demo_w, 0.3,
             "Reviewer-facing homepage for topic discovery, campaign setup, and human approval.",
             size=12, italic=True, color=ORANGE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    preview_placeholder_h = 5.15
    preview_card_x = block_x[0]
    preview_card_y = 16.10
    add_rect(slide, preview_card_x, preview_card_y, demo_w, preview_placeholder_h,
             fill=WHITE, line_color=ORANGE, line_w=2.0, rounded=True)
    insert_logo(
        slide, REVIEWER_PREVIEW,
        preview_card_x + 0.16, preview_card_y + 0.16,
        demo_w - 0.32, preview_placeholder_h - 0.32,
        placeholder_label="reviewer homepage preview",
    )
    stack_card_y = preview_card_y + preview_placeholder_h + 0.22
    stack_card_h = 1.35
    add_rect(slide, block_x[0] + 0.2, stack_card_y, demo_w - 0.4, stack_card_h,
             fill=WHITE, line_color=GRAY_LIGHT, rounded=True)
    add_text(slide, block_x[0] + 0.35, stack_card_y + 0.12, demo_w - 0.7, 0.18,
             "Full-stack tech stack",
             size=11, bold=True, color=BLACK, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    add_text(slide, block_x[0] + 0.35, stack_card_y + 0.34, demo_w - 0.7, 0.22,
             "React + FastAPI + LangGraph",
             size=14, bold=True, color=ORANGE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    add_text(slide, block_x[0] + 0.45, stack_card_y + 0.60, demo_w - 0.9, 0.52,
             [
                 "React powers the reviewer-facing web experience and campaign preview flow.",
                 "FastAPI serves generation and approval APIs, while LangGraph coordinates the multi-agent pipeline and HITL routing.",
             ],
             size=10, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.CENTER)

    # ── BLOCK 2 (middle): Real Campaign Output, 2 generated images ──────────
    _zone_title(block_x[1], demo_title_y, demo_w, "Real Campaign Output")
    img1_path = ROOT / "outputs" / "generated_image_20260417_225704.png"
    img2_path = ROOT / "outputs" / "generated_image_20260412_154909.png"
    img_bbox_w = (demo_w - 0.2) / 2  # two images with 0.2" gap
    img_bbox_y = demo_y + 0.1
    img_bbox_h = demo_h - 0.6  # leave 0.5" for caption below
    cap_y = img_bbox_y + img_bbox_h + 0.03

    for i, (pth, cap) in enumerate([
        (img1_path, "Apr 17 · Pacific Migrant Workers\nStyle: Editorial Portrait"),
        (img2_path, "Apr 12 · Florida — Civil Liberties\nStyle: Cinematic Depth"),
    ]):
        ix = block_x[1] + i * (img_bbox_w + 0.2)
        # Thin frame (slightly larger than image will render)
        insert_logo(slide, pth, ix, img_bbox_y, img_bbox_w, img_bbox_h,
                    placeholder_label=f"image_{i+1}")
        add_text(slide, ix, cap_y, img_bbox_w, 0.5,
                 cap.split("\n"),
                 size=13, color=GRAY_DARK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER)

    # ── BLOCK 3 (right): From Keyword to Campaign — 2-layer pipeline ────────
    _zone_title(block_x[2], demo_title_y, demo_w, "From Keyword to Campaign")

    bx = block_x[2]
    bw = demo_w
    add_text(slide, bx, demo_subtitle_y, bw, 0.30,
             "Two clean swimlanes: discover the signal, then build the campaign.",
             size=12, italic=True, color=ORANGE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)

    def _flow_node(nx, ny, nw, nh, title, tool, detail, *,
                   fill=WHITE, line_color=GRAY_LIGHT,
                   title_color=BLACK, tool_color=ORANGE):
        add_rect(slide, nx, ny, nw, nh, fill=fill, line_color=line_color,
                 line_w=1.6, rounded=True)
        add_text(slide, nx + 0.05, ny + 0.10, nw - 0.10, 0.22, title,
                 size=12, bold=True, color=title_color, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, wrap=False)
        add_text(slide, nx + 0.05, ny + 0.38, nw - 0.10, 0.16, tool,
                 size=9, italic=True, color=tool_color, font=FONT_BODY,
                 align=PP_ALIGN.CENTER, wrap=False)
        add_text(slide, nx + 0.08, ny + 0.64, nw - 0.16, nh - 0.72, detail,
                 size=8.5, color=GRAY_DARK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER, wrap=False)

    def _flow_lane(panel_y, title, subtitle, nodes, note):
        panel_h = 2.15
        add_rect(slide, bx, panel_y, bw, panel_h, fill=WHITE,
                 line_color=GRAY_LIGHT, rounded=True)
        add_text(slide, bx + 0.25, panel_y + 0.18, bw - 0.50, 0.22, title,
                 size=15, bold=True, color=ORANGE, font=FONT_HEAD,
                 wrap=False)
        add_text(slide, bx + 0.25, panel_y + 0.46, bw - 0.50, 0.18, subtitle,
                 size=9.5, italic=True, color=GRAY_DARK, font=FONT_BODY,
                 wrap=False)

        inner_x = bx + 0.25
        node_y = panel_y + 0.78
        node_h = 0.98
        node_gap = 0.18
        node_w = (bw - 0.50 - 3 * node_gap) / 4
        xs = [inner_x + i * (node_w + node_gap) for i in range(4)]

        for i, node in enumerate(nodes):
            title, tool, detail, kind = node
            if kind == "gate":
                _flow_node(xs[i], node_y, node_w, node_h,
                           title, tool, detail,
                           fill=YELLOW, line_color=YELLOW,
                           title_color=BLACK, tool_color=BLACK)
            else:
                _flow_node(xs[i], node_y, node_w, node_h,
                           title, tool, detail)

        mid_y = node_y + node_h / 2
        for i in range(3):
            add_connector_arrow(slide, xs[i] + node_w, mid_y, xs[i + 1], mid_y,
                                color=ORANGE, width_pt=1.5)

        add_text(slide, bx + 0.25, panel_y + panel_h - 0.28, bw - 0.50, 0.16,
                 note, size=8.8, color=GRAY_DARK, font=FONT_BODY,
                 italic=True, wrap=False)

    _flow_lane(
        15.95,
        "1. Trend Discovery",
        "Start from a human topic, then ground the draft in retrieved reporting.",
        [
            ("Topic", "human input", "custom topic or default prompt", "node"),
            ("Search", "Exa + retries", "fresh reporting across trusted domains", "node"),
            ("Extract", "typed facts", "headline, source, and key context", "node"),
            ("HITL Trend", "approve or re-search", "human gate", "gate"),
        ],
        "Fresh articles and structured extraction keep the campaign anchored in real reporting.",
    )

    _flow_lane(
        18.35,
        "2. Campaign Build",
        "Once the topic is approved, audience and image decisions become the key review moments.",
        [
            ("Audience", "voice brief", "who this post is for and how it should sound", "node"),
            ("Prompt + Facts", "writer inputs", "article context plus audience guidance", "node"),
            ("Render + Compose", "Gemini + Pillow", "visual plate, typography, accent bar", "node"),
            ("HITL Image", "approve or regenerate", "human gate", "gate"),
        ],
        "The reviewer can pause, refine, or approve before any sensitive visual ships.",
    )

    add_rect(slide, bx, 20.85, bw, 1.18, fill=GRAY_BG,
             line_color=GRAY_LIGHT, rounded=True)
    add_text(slide, bx + 0.25, 21.03, bw - 0.50, 0.22,
             "Why this layout matters",
             size=13, bold=True, color=BLACK, font=FONT_HEAD, wrap=False)
    add_text(slide, bx + 0.25, 21.35, bw - 0.50, 0.44,
             "The pipeline now reads left-to-right at poster distance: search, extract, draft, review. "
             "Smaller implementation details stay implied instead of fighting for attention.",
             size=9.5, color=GRAY_DARK, font=FONT_BODY)

    # ── TECH STACK HIGHLIGHTS strip (23.35 – 25.15") ────────────────────────
    hi_y = 22.95
    add_text(slide, 1.0, hi_y, W - 2.0, 0.5,
             "TECH STACK — EACH CHOICE SERVES THE MISSION",
             size=18, bold=True, color=ORANGE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    pills = [
        ("EXA.AI",    "Fresh News"),
        ("LANGGRAPH", "Agent Orchestration"),
        ("CHATGPT",   "Audience & Writing"),
        ("GEMINI",    "Campaign Visuals"),
        ("LANGSMITH", "Full Traceability"),
        ("REACT",     "Reviewer UI"),
    ]
    pill_y = hi_y + 0.55
    pill_h = 0.72
    margin = 1.0
    total_w = W - 2 * margin
    n = len(pills)
    gap_px = 0.35
    pill_w = (total_w - (n - 1) * gap_px) / n
    for i, (primary, secondary) in enumerate(pills):
        px = margin + i * (pill_w + gap_px)
        add_stack_chip(slide, px, pill_y, pill_w, pill_h, primary, secondary)

    # ── BOTTOM CARDS (25.3 – 31.8) — Tech Stack | Results | Impact | Future ─
    card_y = 24.80
    card_h = 7.25
    card_w = (W - 2.0 - 3 * 0.4) / 4
    cards_start_x = 1.0

    hood_bullets = [
        "Hallucination-resistant by design — every campaign traces back to an Exa-retrieved news article, not model memory. No fabricated stats, no imagined events.",
        "Prompt engineering as versioned skills — 4 photography presets in visual_styles.md are modular, swappable, Git-tracked; not hardcoded prompts.",
        "Structured output via Pydantic — ChatGPT extracts headline + key fact + source as a typed schema, eliminating free-form parsing bugs.",
        "HITL only at three gates — topic, audience, image; writer and publisher run automatically. Human judgment where it matters; automation where it doesn't.",
        "Typography as post-processing — Pillow + bundled fonts + drop-shadow solve the classic \"AI-generated text looks broken\" problem.",
    ]
    results_bullets = [
        "End-to-end campaign in ~2 min vs. ~45 min manual baseline (writer + designer + editor).",
        "46 campaign images generated over 4-week internship (Mar 16 – Apr 13, 2026).",
        "100% source traceability — every image cites its Exa-retrieved article (n=46 runs).",
        "Relevance-sorted trends cut reviewer article-scan time; Supervisor retries up to 3× on transient errors.",
    ]
    impact_bullets = [
        "Reaches audiences who don't yet see themselves as victims — content tuned around lived-experience cues, not clinical language.",
        "One-person media team — a single reviewer produces what used to need writer + designer + editor.",
        "Modular agent design transfers to other advocacy verticals — labor rights, domestic violence, civic campaigns.",
        "Responsible-AI guardrails — no victim imagery, no identifying details, every post human-approved before publish.",
    ]
    future_bullets = [
        "Quantitative audience-fit evaluation — A/B prompt scoring against a labeled audience-register dataset.",
        "Multi-image / carousel output for richer storytelling.",
        "Agent-level cost & latency dashboard via LangSmith.",
        "Model-agnostic Writer swap (Claude, open-source fallbacks).",
        "Scheduled campaign orchestration via cron triggers.",
    ]
    cards_data = [
        ("Under the Hood", hood_bullets,    20),
        ("Results",        results_bullets, 20),
        ("Impact",         impact_bullets,  20),
        ("Future Work",    future_bullets,  20),
    ]
    for i, (title, bullets, bsize) in enumerate(cards_data):
        cx = cards_start_x + i * (card_w + 0.4)
        add_card(slide, cx, card_y, card_w, card_h, title, bullets, body_size=bsize)

    # ── FOOTER (32.15 – 36) ──────────────────────────────────────────────────
    rule2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.0), Inches(32.30),
                                   Inches(W - 2.0), Inches(0.06))
    set_fill(rule2, ORANGE); no_line(rule2)

    # References (left)
    add_text(slide, 1.0, 32.52, 12.0, 0.42, "References",
             size=18, bold=True, color=BLACK, font=FONT_HEAD)
    add_text(slide, 1.0, 32.94, 12.0, 1.9,
             ["LangGraph — LangChain Inc., 2024.",
               "Exa.ai — Neural Search API docs, 2025.",
               "LangSmith Platform — LangChain observability, 2024.",
               "Capstone Poster Guide 2026, UT iSchool."],
             size=13, color=GRAY_DARK, font=FONT_BODY)

    # AustinAIHUB logo (center)
    insert_logo(slide, austin_prep, W/2 - 4.4, 32.78, 8.8, 1.55,
                placeholder_label="AustinAIHUB logo — drop here")
    add_text(slide, W/2 - 5.4, 34.48, 10.8, 0.4,
             "Sponsored by AustinAIHUB  ·  austinaihub.com",
             size=15, bold=True, color=ORANGE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)

    # Contact (right)
    add_text(slide, W - 10.0, 32.52, 9.0, 0.42, "Contact",
             size=18, bold=True, color=BLACK, font=FONT_HEAD,
             align=PP_ALIGN.RIGHT)
    add_text(slide, W - 10.0, 32.94, 9.0, 1.9,
             ["Kaiwen Lin — AI Engineer Intern",
               "kl38264@utexas.edu",
               "Supervisor: Ali Mirzapour",
               "Team: A. Azami · A. Mohammadi"],
             size=13, color=GRAY_DARK, font=FONT_BODY,
             align=PP_ALIGN.RIGHT)

    out_path = OUT_DIR / "12-Lin-Sp26.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slide size: {prs.slide_width.inches} x {prs.slide_height.inches} inches")


if __name__ == "__main__":
    build()
